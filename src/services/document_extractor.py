from __future__ import annotations

import os
from dataclasses import dataclass
from pathlib import Path
from typing import Optional

import requests
from azure.identity import DefaultAzureCredential


@dataclass
class ExtractedDocument:
    text: str
    kind: str
    used_ocr: bool = False


class DocumentExtractor:
    def __init__(self):
        self.di_endpoint = os.getenv("AZURE_DOCUMENT_INTELLIGENCE_ENDPOINT")
        self.di_api_version = os.getenv("AZURE_DOCUMENT_INTELLIGENCE_API_VERSION", "2023-07-31")

    def extract(self, path: Path, content_type: Optional[str] = None) -> Optional[ExtractedDocument]:
        ext = (path.suffix or "").lower()

        if ext in {".txt", ".md", ".csv", ".json", ".log"}:
            try:
                return ExtractedDocument(text=path.read_text(encoding="utf-8", errors="ignore"), kind=ext.lstrip("."))
            except Exception:
                return None

        if ext == ".pdf":
            text = self._extract_pdf_text(path)
            if text and text.strip():
                return ExtractedDocument(text=text, kind="pdf", used_ocr=False)

            # Scanned PDF: OCR fallback via Azure Document Intelligence if configured
            ocr_text = self._ocr_pdf_with_document_intelligence(path)
            if ocr_text and ocr_text.strip():
                return ExtractedDocument(text=ocr_text, kind="pdf", used_ocr=True)
            return ExtractedDocument(text="", kind="pdf", used_ocr=False)

        if ext == ".docx":
            return self._extract_docx(path)

        if ext == ".pptx":
            return self._extract_pptx(path)

        if ext == ".xlsx":
            return self._extract_xlsx(path)

        return None

    def _extract_pdf_text(self, path: Path) -> str:
        try:
            from pypdf import PdfReader

            reader = PdfReader(str(path))
            parts: list[str] = []
            for page in reader.pages:
                try:
                    parts.append(page.extract_text() or "")
                except Exception:
                    parts.append("")
            return "\n".join(parts)
        except Exception:
            return ""

    def _extract_docx(self, path: Path) -> Optional[ExtractedDocument]:
        try:
            import docx  # python-docx

            d = docx.Document(str(path))
            text = "\n".join([p.text for p in d.paragraphs if p.text])
            return ExtractedDocument(text=text, kind="docx")
        except Exception:
            return None

    def _extract_pptx(self, path: Path) -> Optional[ExtractedDocument]:
        try:
            from pptx import Presentation

            prs = Presentation(str(path))
            parts: list[str] = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        t = (shape.text or "").strip()
                        if t:
                            parts.append(t)
            return ExtractedDocument(text="\n".join(parts), kind="pptx")
        except Exception:
            return None

    def _extract_xlsx(self, path: Path) -> Optional[ExtractedDocument]:
        try:
            from openpyxl import load_workbook

            wb = load_workbook(str(path), read_only=True, data_only=True)
            parts: list[str] = []
            for sheet in wb.worksheets:
                parts.append(f"[Sheet: {sheet.title}]")
                for row in sheet.iter_rows(values_only=True):
                    vals = ["" if v is None else str(v) for v in row]
                    line = "\t".join([v for v in vals if v != ""]) 
                    if line.strip():
                        parts.append(line)
            return ExtractedDocument(text="\n".join(parts), kind="xlsx")
        except Exception:
            return None

    def _ocr_pdf_with_document_intelligence(self, path: Path) -> str:
        if not self.di_endpoint:
            return ""

        try:
            credential = DefaultAzureCredential()
            token = credential.get_token("https://cognitiveservices.azure.com/.default").token
        except Exception:
            return ""

        url = (
            f"{self.di_endpoint.rstrip('/')}/documentintelligence/documentModels/prebuilt-read:analyze"
            f"?api-version={self.di_api_version}"
        )

        try:
            with open(path, "rb") as f:
                resp = requests.post(
                    url,
                    headers={
                        "Authorization": f"Bearer {token}",
                        "Content-Type": "application/pdf",
                    },
                    data=f,
                    timeout=120,
                )
            resp.raise_for_status()
            op_loc = resp.headers.get("operation-location")
            if not op_loc:
                return ""

            # Poll
            for _ in range(60):
                poll = requests.get(op_loc, headers={"Authorization": f"Bearer {token}"}, timeout=60)
                if poll.status_code >= 400:
                    return ""
                data = poll.json()
                status = (data.get("status") or "").lower()
                if status in {"succeeded", "failed"}:
                    if status != "succeeded":
                        return ""
                    break
                import time

                time.sleep(1)

            analyze = data.get("analyzeResult") or {}
            content = analyze.get("content")
            if isinstance(content, str):
                return content
            return ""
        except Exception:
            return ""
